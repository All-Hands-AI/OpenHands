"""File reader skills for the OpenHands agent.

This module provides various functions to parse and extract content from different file types,
including PDF, DOCX, LaTeX, audio, image, video, and PowerPoint files. It utilizes different
libraries and APIs to process these files and output their content or descriptions.

Functions:
    parse_pdf(file_path: str) -> None: Parse and print content of a PDF file.
    parse_docx(file_path: str) -> None: Parse and print content of a DOCX file.
    parse_latex(file_path: str) -> None: Parse and print content of a LaTeX file.
    parse_audio(file_path: str, model: str = 'whisper-1') -> None: Transcribe and print content of an audio file.
    parse_image(file_path: str, task: str = 'Describe this image as detail as possible.') -> None: Analyze and print description of an image file.
    parse_video(file_path: str, task: str = 'Describe this image as detail as possible.', frame_interval: int = 30) -> None: Analyze and print description of video frames.
    parse_pptx(file_path: str) -> None: Parse and print content of a PowerPoint file.

Note:
    Some functions (parse_audio, parse_video, parse_image) require OpenAI API credentials
    and are only available if the necessary environment variables are set.
"""

import base64
from typing import Any

import docx
import PyPDF2
from pptx import Presentation
from pylatexenc.latex2text import LatexNodes2Text

from openhands.runtime.plugins.agent_skills.utils.config import (
    _get_max_token,
    _get_openai_api_key,
    _get_openai_base_url,
    _get_openai_client,
    _get_openai_model,
)


def parse_pdf(file_path: str) -> None:
    """Parses the content of a PDF file and prints it.

    Args:
        file_path: str: The path to the file to open.
    """
    print(f'[Reading PDF file from {file_path}]')
    content = PyPDF2.PdfReader(file_path)
    text = ''
    for page_idx in range(len(content.pages)):
        text += (
            f'@@ Page {page_idx + 1} @@\n'
            + content.pages[page_idx].extract_text()
            + '\n\n'
        )
    print(text.strip())


def parse_docx(file_path: str) -> None:
    """Parses the content of a DOCX file and prints it.

    Args:
        file_path: str: The path to the file to open.
    """
    print(f'[Reading DOCX file from {file_path}]')
    content = docx.Document(file_path)
    text = ''
    for i, para in enumerate(content.paragraphs):
        text += f'@@ Page {i + 1} @@\n' + para.text + '\n\n'
    print(text)


def parse_latex(file_path: str) -> None:
    """Parses the content of a LaTex file and prints it.

    Args:
        file_path: str: The path to the file to open.
    """
    print(f'[Reading LaTex file from {file_path}]')
    with open(file_path) as f:
        data = f.read()
    text = LatexNodes2Text().latex_to_text(data)
    print(text.strip())


def _base64_img(file_path: str) -> str:
    with open(file_path, 'rb') as image_file:
        encoded_image = base64.b64encode(image_file.read()).decode('utf-8')
    return encoded_image


def _base64_video(file_path: str, frame_interval: int = 10) -> list[str]:
    import cv2

    video = cv2.VideoCapture(file_path)
    base64_frames = []
    frame_count = 0
    while video.isOpened():
        success, frame = video.read()
        if not success:
            break
        if frame_count % frame_interval == 0:
            _, buffer = cv2.imencode('.jpg', frame)
            base64_frames.append(base64.b64encode(buffer).decode('utf-8'))
        frame_count += 1
    video.release()
    return base64_frames


def _prepare_image_messages(task: str, base64_image: str) -> list[dict[str, Any]]:
    return [
        {
            'role': 'user',
            'content': [
                {'type': 'text', 'text': task},
                {
                    'type': 'image_url',
                    'image_url': {'url': f'data:image/jpeg;base64,{base64_image}'},
                },
            ],
        }
    ]


def parse_audio(file_path: str, model: str = 'whisper-1') -> None:
    """Parses the content of an audio file and prints it.

    Args:
        file_path: str: The path to the audio file to transcribe.
        model: str: The audio model to use for transcription. Defaults to 'whisper-1'.
    """
    print(f'[Transcribing audio file from {file_path}]')
    try:
        # TODO: record the COST of the API call
        with open(file_path, 'rb') as audio_file:
            transcript = _get_openai_client().audio.translations.create(
                model=model, file=audio_file
            )
        print(transcript.text)

    except Exception as e:
        print(f'Error transcribing audio file: {e}')


def parse_image(
    file_path: str, task: str = 'Describe this image as detail as possible.'
) -> None:
    """Parses the content of an image file and prints the description.

    Args:
        file_path: str: The path to the file to open.
        task: str: The task description for the API call. Defaults to 'Describe this image as detail as possible.'.
    """
    print(f'[Reading image file from {file_path}]')
    # TODO: record the COST of the API call
    try:
        base64_image = _base64_img(file_path)
        response = _get_openai_client().chat.completions.create(
            model=_get_openai_model(),
            messages=_prepare_image_messages(task, base64_image),
            max_tokens=_get_max_token(),
        )
        content = response.choices[0].message.content
        print(content)

    except Exception as error:
        print(f'Error with the request: {error}')


def parse_video(
    file_path: str,
    task: str = 'Describe this image as detail as possible.',
    frame_interval: int = 30,
) -> None:
    """Parses the content of an image file and prints the description.

    Args:
        file_path: str: The path to the video file to open.
        task: str: The task description for the API call. Defaults to 'Describe this image as detail as possible.'.
        frame_interval: int: The interval between frames to analyze. Defaults to 30.

    """
    print(
        f'[Processing video file from {file_path} with frame interval {frame_interval}]'
    )

    task = task or 'This is one frame from a video, please summarize this frame.'
    base64_frames = _base64_video(file_path)
    selected_frames = base64_frames[::frame_interval]

    if len(selected_frames) > 30:
        new_interval = len(base64_frames) // 30
        selected_frames = base64_frames[::new_interval]

    print(f'Totally {len(selected_frames)} would be analyze...\n')

    idx = 0
    for base64_frame in selected_frames:
        idx += 1
        print(f'Process the {file_path}, current No. {idx * frame_interval} frame...')
        # TODO: record the COST of the API call
        try:
            response = _get_openai_client().chat.completions.create(
                model=_get_openai_model(),
                messages=_prepare_image_messages(task, base64_frame),
                max_tokens=_get_max_token(),
            )

            content = response.choices[0].message.content
            current_frame_content = f"Frame {idx}'s content: {content}\n"
            print(current_frame_content)

        except Exception as error:
            print(f'Error with the request: {error}')


def parse_pptx(file_path: str) -> None:
    """Parses the content of a pptx file and prints it.

    Args:
        file_path: str: The path to the file to open.
    """
    print(f'[Reading PowerPoint file from {file_path}]')
    try:
        pres = Presentation(str(file_path))
        text = []
        for slide_idx, slide in enumerate(pres.slides):
            text.append(f'@@ Slide {slide_idx + 1} @@')
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text.append(shape.text)
        print('\n'.join(text))

    except Exception as e:
        print(f'Error reading PowerPoint file: {e}')


__all__ = [
    'parse_pdf',
    'parse_docx',
    'parse_latex',
    'parse_pptx',
]

# This is called from OpenHands's side
# If SANDBOX_ENV_OPENAI_API_KEY is set, we will be able to use these tools in the sandbox environment
if _get_openai_api_key() and _get_openai_base_url():
    __all__ += ['parse_audio', 'parse_video', 'parse_image']
