import base64
import json
import subprocess
import zipfile
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Any, Union

import charset_normalizer
import cv2
import docx

# import markdown
import openpyxl
import pandas as pd
import PyPDF2
import requests
import yaml
from bs4 import BeautifulSoup
from openai import OpenAI
from pptx import Presentation
from pylatexenc.latex2text import LatexNodes2Text

from opendevin.core.config import config
from opendevin.core.logger import opendevin_logger as logger

# TODO: Find way to directly get the API key from ConfigType.LLM_API_KEY or change it with litellm.
OPENAI_API_KEY = config.llm.api_key


class Reader(ABC):
    """
    @Desc: Implementation to support reading 41 multimodal files efficiently.
    @Ref: https://github.com/metauto-ai/GPTSwarm/blob/main/swarm/environment/tools/reader/readers.py
    """

    @abstractmethod
    def parse(self, file_path: Path) -> str:
        """To be overriden by the descendant class"""
        pass


class TXTReader(Reader):
    def parse(self, file_path: Path) -> str:
        content = charset_normalizer.from_path(file_path).best()
        logger.info(
            f"Reading TXT file from {file_path} using encoding '{content.encoding}.'"
        )
        return str(content)


class PDFReader(Reader):
    def parse(self, file_path: Path) -> str:
        logger.info(f'Reading PDF file from {file_path}.')
        content = PyPDF2.PdfReader(file_path)
        text = ''
        for page_idx in range(len(content.pages)):
            text += f'Page {page_idx + 1}\n' + content.pages[page_idx].extract_text()
        return text


class DOCXReader(Reader):
    def parse(self, file_path: Path) -> str:
        logger.info(f'Reading DOCX file from {file_path}.')
        content = docx.Document(str(file_path))
        text = ''
        for i, para in enumerate(content.paragraphs):
            text += f'Page {i + 1}:\n' + para.text
        return text


class JSONReader(Reader):
    def parse(self, file_path: Path) -> str:
        logger.info(f'Reading JSON file from {file_path}.')

        with open(file_path, 'r') as f:
            data = json.load(f)
            text = str(data)
        return text


class JSONLReader(Reader):
    def parse(self, file_path: Path) -> str:
        logger.info(f'Reading JSON Lines file from {file_path}.')
        with open(file_path, 'r') as f:
            lines = [json.loads(line) for line in f]
            text = '\n'.join([str(line) for line in lines])
        return text


class XMLReader(Reader):
    def parse(self, file_path: Path) -> str:
        logger.info(f'Reading XML file from {file_path}.')
        with open(file_path, 'r') as f:
            data = BeautifulSoup(f, 'xml')
            text = data.get_text()
        return text


class YAMLReader(Reader):
    def parse(self, file_path: Path, return_str=True) -> Union[str, Any]:
        logger.info(f'Reading YAML file from {file_path}.')
        with open(file_path, 'r') as f:
            data = yaml.load(f, Loader=yaml.FullLoader)
            text = str(data)
        if return_str:
            return text
        else:
            return data


class HTMLReader(Reader):
    def parse(self, file_path: Path) -> str:
        logger.info(f'Reading HTML file from {file_path}.')
        with open(file_path, 'r') as f:
            data = BeautifulSoup(f, 'html.parser')
            text = data.get_text()
        return text


# class MarkdownReader(Reader):
#     def parse(self, file_path: Path) -> str:
#         logger.info(f'Reading Markdown file from {file_path}.')
#         with open(file_path, 'r') as f:
#             data = markdown.markdown(f.read())
#             text = ''.join(BeautifulSoup(data, 'html.parser').findAll(string=True))
#         return text


class LaTexReader(Reader):
    def parse(self, file_path: Path) -> str:
        logger.info(f'Reading LaTex file from {file_path}.')
        with open(file_path, 'r') as f:
            data = f.read()
        text = LatexNodes2Text().latex_to_text(data)
        return text


class AudioReader(Reader):
    @staticmethod
    def parse(file_path: Path) -> str:
        api_key = OPENAI_API_KEY
        logger.info(f'Transcribing audio file from {file_path}.')
        client = OpenAI(api_key=api_key)
        try:
            # TODO: record the COST of the API call
            client = OpenAI()
            with open(file_path, 'rb') as audio_file:
                transcript = client.audio.translations.create(
                    model='whisper-1', file=audio_file
                )
            return transcript.text

        except Exception as e:
            logger.info(f'Error transcribing audio file: {e}')
            return 'Error transcribing audio file.'


class PPTXReader(Reader):
    def parse(self, file_path: Path) -> str:
        logger.info(f'Reading PowerPoint file from {file_path}.')
        try:
            pres = Presentation(str(file_path))
            text = []
            for slide_idx, slide in enumerate(pres.slides):
                text.append(f'Slide {slide_idx + 1}:\n')
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text.append(shape.text)
            return '\n'.join(text)

        except Exception as e:
            logger.info(f'Error reading PowerPoint file: {e}')
            return 'Error reading PowerPoint file.'


class ExcelReader(Reader):
    def parse(self, file_path: Path) -> str:
        logger.info(f'Reading Excel file from {file_path}.')
        try:
            excel_data = pd.read_excel(file_path, sheet_name=None)

            all_sheets_text = []
            for sheet_name, data in excel_data.items():
                all_sheets_text.append(
                    f'Sheet Name: {sheet_name}\n{data.to_string()}\n'
                )
            return '\n'.join(all_sheets_text)

        except Exception as e:
            logger.info(f'Error reading Excel file: {e}')
            return 'Error reading Excel file.'


class XLSXReader(Reader):
    def parse(self, file_path: Path) -> str:
        logger.info(f'Reading XLSX file from {file_path}.')
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        text = ''

        for sheet in workbook:
            text += f'\nSheet: {sheet.title}\n'
            for row in sheet.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else '' for cell in row]
                text += '\t'.join(row_data) + '\n'
        return text


class ZipReader(Reader):
    def parse(self, file_path: Path) -> str:
        logger.info(f'Reading ZIP file from {file_path}.')
        file_content = ''
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            extract_dir = str(file_path)[:-4] + '/'
            zip_ref.extractall(Path(extract_dir))
            reader = FileReader()
            for file_name in zip_ref.namelist():
                file_content += f'File {file_name}:\n"{reader.read_file(Path(extract_dir + file_name))}"\n'
        return file_content


class PythonReader(Reader):
    def parse(self, file_path: Path) -> str:
        logger.info(f'Executing and reading Python file from {file_path}.')
        execution_result = ''
        try:
            completed_process = subprocess.run(
                ['python', file_path], capture_output=True, text=True, check=True
            )
            execution_result = 'Execution information:\n' + completed_process.stdout
        except subprocess.CalledProcessError as e:
            execution_result = 'Error:\n' + e.stderr
            return execution_result
        except Exception as e:
            logger.info(f'Error executing Python file: {e}')

        try:
            with open(file_path, 'r') as file:
                file_content = '\nFile Content:\n' + file.read()
        except Exception as e:
            logger.info(f'Error reading Python file: {e}')
        return file_content + '\n' + execution_result


class IMGReader(Reader):
    def base64_img(self, file_path: Path) -> str:
        import base64

        with open(file_path, 'rb') as image_file:
            encoded_image = base64.b64encode(image_file.read()).decode('utf-8')
        return encoded_image

    def prepare_api_call(
        self, task: str, base64_frame: str, model='gpt-4o-2024-05-13', max_token=500
    ) -> dict:
        return {
            'model': model,
            'messages': [
                {
                    'role': 'user',
                    'content': [
                        {'type': 'text', 'text': task},
                        {
                            'type': 'image_url',
                            'image_url': {
                                'url': f'data:image/jpeg;base64,{base64_frame}'
                            },
                        },
                    ],
                }
            ],
            'max_tokens': max_token,
        }

    def get_headers(self) -> dict:
        return {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {OPENAI_API_KEY}',
        }

    def parse(
        self, file_path: Path, task: str = 'Describe this image as detail as possible.'
    ) -> str:
        logger.info(f'Reading image file from {file_path}.')
        # TODO: record the COST of the API call
        try:
            openai_proxy: str = 'https://api.openai.com/v1/chat/completions'
            base64_image = self.base64_img(Path(file_path))
            api_call = self.prepare_api_call(task, base64_image)
            response = requests.post(
                openai_proxy, headers=self.get_headers(), json=api_call
            )
            out = response.json()
            content = out['choices'][0]['message']['content']
            return content

        except Exception as error:
            logger.error(f'Error with the request: {error}')
            raise


class VideoReader(Reader):
    def base64_video(self, file_path: Path, frame_interval: int = 10) -> list:
        video = cv2.VideoCapture(str(file_path))
        base64_frames = []
        frame_count = 0
        while video.isOpened():
            success, frame = video.read()
            if not success:
                break
            if frame_count % frame_interval == 0:
                _, buffer = cv2.imencode('.jpg', frame)
                base64_frames.append(base64.b64encode(buffer).decode('utf-8'))
            frame_count += 1
        video.release()
        return base64_frames

    def prepare_api_call(
        self, task: str, base64_frame: str, model='gpt-4o-2024-05-13', max_token=500
    ) -> dict:
        return {
            'model': model,
            'messages': [
                {
                    'role': 'user',
                    'content': [
                        {'type': 'text', 'text': task},
                        {
                            'type': 'image_url',
                            'image_url': {
                                'url': f'data:image/jpeg;base64,{base64_frame}'
                            },
                        },
                    ],
                }
            ],
            'max_tokens': max_token,
        }

    def get_headers(self) -> dict:
        return {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {OPENAI_API_KEY}',
        }

    def parse(
        self,
        file_path: Path,
        task: str = 'Describe this image as detail as possible.',
        frame_interval: int = 30,
        used_audio: bool = True,
    ) -> str:
        logger.info(
            f'Processing video file from {file_path} with frame interval {frame_interval}.'
        )

        video_summary = ''
        idx = 0
        task = task or 'This is one frame from a video, please summarize this frame.'
        base64_frames = self.base64_video(Path(file_path))
        selected_frames = base64_frames[::frame_interval]

        if len(selected_frames) > 30:
            new_interval = len(base64_frames) // 30
            selected_frames = base64_frames[::new_interval]

        logger.info(f'Totally {len(selected_frames)} would be analyze...')

        idx = 0
        for base64_frame in selected_frames:
            idx += 1
            logger.info(
                f'Process the {file_path}, current No. {idx * frame_interval} frame...'
            )
            # TODO: record the COST of the API call
            api_call = self.prepare_api_call(task, base64_frame)
            try:
                openai_proxy: str = 'https://api.openai.com/v1/chat/completions'
                response = requests.post(
                    openai_proxy, headers=self.get_headers(), json=api_call
                )
                content = response.json()['choices'][0]['message']['content']
                current_frame_content = f"Frame {idx}'s content: {content}\n"
                video_summary += current_frame_content
                logger.info(current_frame_content)

            except Exception as error:
                logger.error(f'Error with the request: {error}')
                raise

        logger.info(f'video summary: {video_summary}')
        return video_summary


# Support 41 kinds of files.
READER_MAP = {
    '.png': IMGReader(),
    '.jpg': IMGReader(),
    '.jpeg': IMGReader(),
    '.gif': IMGReader(),
    '.bmp': IMGReader(),
    '.tiff': IMGReader(),
    '.tif': IMGReader(),
    '.webp': IMGReader(),
    '.mp3': AudioReader(),
    '.m4a': AudioReader(),
    '.wav': AudioReader(),
    '.MOV': VideoReader(),
    '.mp4': VideoReader(),
    '.mov': VideoReader(),
    '.avi': VideoReader(),
    '.mpg': VideoReader(),
    '.mpeg': VideoReader(),
    '.wmv': VideoReader(),
    '.flv': VideoReader(),
    '.webm': VideoReader(),
    '.zip': ZipReader(),
    '.pptx': PPTXReader(),
    '.xlsx': ExcelReader(),
    '.xls': ExcelReader(),
    '.txt': TXTReader(),
    '.csv': TXTReader(),
    '.pdf': PDFReader(),
    '.docx': DOCXReader(),
    '.json': JSONReader(),
    '.jsonld': JSONReader(),
    '.jsonl': JSONLReader(),
    '.xml': XMLReader(),
    '.yaml': YAMLReader(),
    '.yml': YAMLReader(),
    '.html': HTMLReader(),
    '.htm': HTMLReader(),
    '.xhtml': HTMLReader(),
    # '.md': MarkdownReader(),
    # '.markdown': MarkdownReader(),
    '.tex': LaTexReader(),
    '.py': PythonReader(),
    '.pdb': TXTReader(),
}


class FileReader:
    def set_reader(self, suffix) -> None:
        self.reader = READER_MAP[suffix]
        logger.info(f'Setting Reader to {type(self.reader).__name__}')

    def read_file(self, file_path: Path, task='describe the file') -> str:
        suffix = file_path.suffix
        self.set_reader(suffix)
        if isinstance(self.reader, IMGReader) or isinstance(self.reader, VideoReader):
            file_content = self.reader.parse(file_path, task)
        else:
            file_content = self.reader.parse(file_path)
        logger.info(f'Reading file {file_path} using {type(self.reader).__name__}')
        return file_content


class GeneralReader:
    def __init__(self):
        self.file_reader = FileReader()
        self.name = 'General File Reader'
        self.description = """A general file reader support to multimodal files."""

    def read(self, task, file):
        files_content = ''
        file_content = self.file_reader.read_file(file, task)
        suffix = file.split('.')[-1]

        if suffix in ['py', 'java', 'cpp', 'c', 'js', 'css', 'html', 'htm', 'xml']:
            files_content += f'\nThe {suffix} file contains:\n---\n{file_content[0]}'
            if file_content[1] != '':
                files_content += f'\nExecution result:\n{file_content[1]}'
            if file_content[2] != '':
                files_content += f'\nExecution error message:\n{file_content[2]}'
            files_content += '\n---'

        elif suffix in [
            'txt',
            'jsonl',
            'csv',
            'json',
            'jsonld',
            'jsonl',
            'yaml',
            'yml',
            'xlsx',
            'xls',
            'jpg',
            'png',
            'jpeg',
            'gif',
            'bmp',
            'mp3',
            'wav',
            'ogg',
            'mp4',
            'avi',
            'mkv',
            'mov',
            'pdf',
            'doc',
            'docx',
            'ppt',
            'pptx',
            'md',
            'markdown',
            'tex',
            'zip',
            'tar',
            'gz',
            '7z',
            'rar',
        ]:
            files_content += f'\nThe {suffix} file contains:\n---\n{file_content}\n---'

        return files_content
